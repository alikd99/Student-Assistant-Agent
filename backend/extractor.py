"""
extractor.py — unified document text extractor
Supports: PDF (native + scanned/OCR), PPTX, DOCX
Each page/slide is returned as a separate ExtractedPage so chunk metadata
carries location info (filename, page/slide number).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class ExtractedPage:
    """One page (PDF) or one slide (PPTX) or one paragraph-block (DOCX)."""
    filename: str
    page_num: int          # 1-indexed
    text: str
    source_type: Literal["pdf", "pptx", "docx"]
    ocr_used: bool = False
    warning: str | None = None   # e.g. "no text found — image-only slide"


@dataclass
class ExtractionResult:
    filename: str
    pages: list[ExtractedPage] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())

    @property
    def total_chars(self) -> int:
        return sum(len(p.text) for p in self.pages)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Minimum chars-per-page before we assume the page is scanned and needs OCR
_OCR_THRESHOLD = 50


def _needs_ocr(text: str) -> bool:
    return len(text.strip()) < _OCR_THRESHOLD


def _ocr_image(image) -> str:
    """Run Tesseract OCR on a PIL Image. Returns extracted text."""
    try:
        import pytesseract
    except ImportError:
        logger.warning("pytesseract not installed — OCR skipped")
        return ""

    try:
        # Support Arabic + English; oem 3 = best LSTM engine, psm 3 = auto layout
        custom_config = r"--oem 3 --psm 3 -l ara+eng"
        text = pytesseract.image_to_string(image, config=custom_config)
        return text.strip()
    except pytesseract.TesseractNotFoundError:
        logger.warning("Tesseract not found — install from https://github.com/UB-Mannheim/tesseract/wiki")
        return ""
    except Exception as exc:
        logger.warning("OCR failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# PDF extractor
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path, filename: str) -> ExtractionResult:
    try:
        import pdfplumber
    except ImportError:
        return ExtractionResult(filename=filename, errors=["pdfplumber not installed"])

    result = ExtractionResult(filename=filename)

    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                native_text: str = page.extract_text() or ""

                if _needs_ocr(native_text):
                    # Render page to image and run OCR
                    try:
                        pil_image = page.to_image(resolution=200).original
                        ocr_text = _ocr_image(pil_image)
                        if ocr_text.strip():
                            result.pages.append(
                                ExtractedPage(
                                    filename=filename,
                                    page_num=i,
                                    text=ocr_text,
                                    source_type="pdf",
                                    ocr_used=True,
                                )
                            )
                        else:
                            result.pages.append(
                                ExtractedPage(
                                    filename=filename,
                                    page_num=i,
                                    text="",
                                    source_type="pdf",
                                    ocr_used=True,
                                    warning="no text found after OCR — possibly blank or image-only page",
                                )
                            )
                    except Exception as exc:
                        result.pages.append(
                            ExtractedPage(
                                filename=filename,
                                page_num=i,
                                text=native_text,
                                source_type="pdf",
                                warning=f"OCR attempt failed: {exc}",
                            )
                        )
                else:
                    result.pages.append(
                        ExtractedPage(
                            filename=filename,
                            page_num=i,
                            text=native_text,
                            source_type="pdf",
                        )
                    )
    except Exception as exc:
        result.errors.append(f"Failed to open PDF: {exc}")

    return result


# ---------------------------------------------------------------------------
# PPTX extractor
# ---------------------------------------------------------------------------

def _extract_pptx(path: Path, filename: str) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(filename=filename, errors=["python-pptx not installed"])

    result = ExtractionResult(filename=filename)

    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []

            for shape in slide.shapes:
                # Text frames (handles RTL Arabic runs correctly)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = "".join(run.text for run in para.runs).strip()
                        if para_text:
                            texts.append(para_text)

                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)

            combined = "\n".join(texts)

            if not combined.strip():
                # Slide has no extractable text — likely image-only
                result.pages.append(
                    ExtractedPage(
                        filename=filename,
                        page_num=i,
                        text="",
                        source_type="pptx",
                        warning="no text found — slide may be image-only",
                    )
                )
            else:
                result.pages.append(
                    ExtractedPage(
                        filename=filename,
                        page_num=i,
                        text=combined,
                        source_type="pptx",
                    )
                )
    except Exception as exc:
        result.errors.append(f"Failed to open PPTX: {exc}")

    return result


# ---------------------------------------------------------------------------
# DOCX extractor
# ---------------------------------------------------------------------------

def _extract_docx(path: Path, filename: str) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(filename=filename, errors=["python-docx not installed"])

    result = ExtractionResult(filename=filename)

    try:
        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        if not paragraphs:
            result.errors.append("No text found in DOCX — document may be empty")
            return result

        # DOCX has no natural page boundaries — group into blocks of 20 paragraphs
        # so each ExtractedPage carries a manageable amount of text
        block_size = 20
        for block_num, start in enumerate(range(0, len(paragraphs), block_size), start=1):
            block = paragraphs[start : start + block_size]
            result.pages.append(
                ExtractedPage(
                    filename=filename,
                    page_num=block_num,
                    text="\n".join(block),
                    source_type="docx",
                )
            )
    except Exception as exc:
        result.errors.append(f"Failed to open DOCX: {exc}")

    return result


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

_SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx"}


def extract(file_path: str | Path, filename: str | None = None) -> ExtractionResult:
    """
    Extract text from a PDF, PPTX, or DOCX file.

    Args:
        file_path: Path to the file on disk.
        filename:  Display name (defaults to the file's stem).

    Returns:
        ExtractionResult with a list of ExtractedPage objects.
    """
    path = Path(file_path)
    display_name = filename or path.name
    ext = path.suffix.lower()

    if ext not in _SUPPORTED_EXTENSIONS:
        return ExtractionResult(
            filename=display_name,
            errors=[f"Unsupported file type '{ext}'. Supported: {_SUPPORTED_EXTENSIONS}"],
        )

    if not path.exists():
        return ExtractionResult(
            filename=display_name,
            errors=[f"File not found: {path}"],
        )

    logger.info("Extracting '%s' (type=%s)", display_name, ext)

    if ext == ".pdf":
        return _extract_pdf(path, display_name)
    elif ext == ".pptx":
        return _extract_pptx(path, display_name)
    elif ext == ".docx":
        return _extract_docx(path, display_name)


# ---------------------------------------------------------------------------
# Quick CLI test  —  python extractor.py <path-to-file>
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    if len(sys.argv) < 2:
        print("Usage: python extractor.py <file.pdf|file.pptx|file.docx>")
        sys.exit(1)

    res = extract(sys.argv[1])
    print(f"\nFile     : {res.filename}")
    print(f"Pages    : {len(res.pages)}")
    print(f"Total chars: {res.total_chars:,}")

    if res.errors:
        print("\nErrors:")
        for e in res.errors:
            print(f"  ✗ {e}")

    print("\n--- Pages ---")
    for page in res.pages:
        ocr_tag = " [OCR]" if page.ocr_used else ""
        warn_tag = f" ⚠ {page.warning}" if page.warning else ""
        preview = page.text[:120].replace("\n", " ") if page.text else "(empty)"
        print(f"  [{page.source_type} p{page.page_num}]{ocr_tag}{warn_tag}  {preview}")
